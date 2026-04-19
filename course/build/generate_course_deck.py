"""Generate ONE course slide deck from all instructor configs, Google-Slides friendly.

Design principles:
- Larger fonts (body 18-22pt, titles 36pt+)
- Fewer items per slide; split into multiple slides if needed
- Generous margins (1" all around)
- Simple layouts; minimal shape overlap
"""
from __future__ import annotations
import json
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor


_HERE = Path(__file__).resolve().parent  # course/build/
_COURSE = _HERE.parent                    # course/
CONFIGS_DIR = _COURSE / "instructor" / "configs"
OUT = _COURSE / "slides" / "course_deck.pptx"
IMG = _COURSE / "shared" / "images"

NAVY = RGBColor(0x0B, 0x2B, 0x5C)
RED = RGBColor(0xE0, 0x29, 0x29)
YELLOW = RGBColor(0xFF, 0xCB, 0x05)
GREEN = RGBColor(0x2E, 0x8B, 0x3F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xF5, 0xF6, 0xFA)

SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN = 0.8

LESSON = {
    "01_intro": {"title": "Intro to Tableau", "subtitle": "Interface tour, dimensions vs measures", "image": "pikachu.jpg", "duration": "45 min", "day": 1},
    "02_basic_visuals": {"title": "Basic Visuals", "subtitle": "Bar, line, scatter, pie, treemap, heatmap", "image": "bulbasaur.jpg", "duration": "60 min", "day": 1},
    "03_advanced_visuals": {"title": "Advanced Visuals", "subtitle": "Time series and multi-encoding marks", "image": "charizard.jpg", "duration": "60 min", "day": 1},
    "04_joins_intro": {"title": "Joins & Data Modeling Intro", "subtitle": "Join battles with venue", "image": "lucario.jpg", "duration": "60 min", "day": 1},
    "05_joins_advanced": {"title": "Joins + Advanced Visualizations", "subtitle": "Three-table join", "image": "greninja.jpg", "duration": "60 min", "day": 1},
    "06_story": {"title": "Day 1 Story Recap", "subtitle": "Tableau Stories", "image": "mewtwo.jpg", "duration": "45 min", "day": 1},
    "07_warmup": {"title": "Day 2 Warm-up", "subtitle": "Trainer exploration", "image": "magikarp.jpg", "duration": "30 min", "day": 2},
    "08_groups_sets": {"title": "Groups & Sets", "subtitle": "Custom segmentation", "image": "snorlax.jpg", "duration": "60 min", "day": 2},
    "09_calculations": {"title": "Calculations", "subtitle": "Conditional logic, aggregates, math", "image": "gengar.jpg", "duration": "60 min", "day": 2},
    "10_mapping": {"title": "Mapping", "subtitle": "Geographic visualization", "image": "dragonite.jpg", "duration": "45 min", "day": 2},
    "11_dashboards": {"title": "Dashboards", "subtitle": "Combine views, add interactivity", "image": "pikachu.jpg", "duration": "60 min", "day": 2},
    "12_lod": {"title": "LOD Calculations", "subtitle": "FIXED, INCLUDE, EXCLUDE", "image": "mewtwo.jpg", "duration": "60 min", "day": 2},
    "13_master_reference": {"title": "Master Reference", "subtitle": "Instructor apex workbook", "image": "charizard.jpg", "duration": "office hours", "day": 2},
}


def blank_slide(prs, bg=WHITE):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bgr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(SLIDE_W), Inches(SLIDE_H))
    bgr.fill.solid(); bgr.fill.fore_color.rgb = bg; bgr.line.fill.background()
    sp = s.shapes._spTree
    sp.remove(bgr._element); sp.insert(2, bgr._element)
    return s


def text(slide, left, top, width, height, content, *, size=20, bold=False, color=DARK, align="left", anchor="top"):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}[anchor]
    lines = content.split("\n") if isinstance(content, str) else content
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
        run = p.add_run(); run.text = line
        run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = color


def bullets(slide, left, top, width, height, items, *, size=20, color=DARK, spacing=8):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    for i, b in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(spacing)
        p.alignment = PP_ALIGN.LEFT
        # Sub-bullet if starts with 2+ spaces
        is_sub = b.startswith("  ")
        prefix = "◦  " if is_sub else "•  "
        run = p.add_run()
        run.text = prefix + b.lstrip()
        run.font.size = Pt(size - 2 if is_sub else size)
        run.font.color.rgb = color


def image(slide, left, top, height, name):
    p = IMG / name
    if p.exists():
        slide.shapes.add_picture(str(p), Inches(left), Inches(top), height=Inches(height))


def header_bar(slide, color, title_text):
    """Top bar, 1.2" tall, with a big centered title."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(SLIDE_W), Inches(1.2))
    bar.fill.solid(); bar.fill.fore_color.rgb = color; bar.line.fill.background()
    text(slide, MARGIN, 0.25, SLIDE_W - 2 * MARGIN, 0.8,
         title_text, size=28, bold=True, color=WHITE, anchor="middle")


# ---------------- Frames ----------------

def title_slide(prs):
    s = blank_slide(prs, NAVY)
    text(s, MARGIN, 1.6, SLIDE_W - 2 * MARGIN, 1.6,
         "Pokémon Tableau Course", size=64, bold=True, color=WHITE, align="center")
    text(s, MARGIN, 3.4, SLIDE_W - 2 * MARGIN, 0.7,
         "13 lessons · 2 days · one Pokémon battle dataset",
         size=24, color=YELLOW, align="center")
    text(s, MARGIN, 4.3, SLIDE_W - 2 * MARGIN, 0.6,
         "One deck for instructors and students — use side by side with the lesson plans.",
         size=16, color=WHITE, align="center")
    image(s, SLIDE_W / 2 - 0.9, 5.2, 1.8, "pikachu.jpg")


def how_to_use(prs):
    s = blank_slide(prs)
    header_bar(s, NAVY, "How this course is organized")
    bullets(s, MARGIN, 1.7, SLIDE_W - 2 * MARGIN, 5.2, [
        "📘 Instructor LP — your script, sheet by sheet  (course/shared/lesson_plans/)",
        "📊 Instructor workbook — pre-built demo  (course/instructor/workbooks/)",
        "🎓 Student blank workbook — students start here  (course/student/workbooks_blank/)",
        "✅ Student answer key — compare after attempting  (course/student/workbooks/)",
        "🗂️  Raw data — every CSV the lessons use  (course/shared/raw_data/)",
        "🖼️  This deck — open per lesson alongside the workbook",
    ], size=18, spacing=14)


def per_lesson_layout(prs, day):
    s = blank_slide(prs)
    header_bar(s, NAVY, "Per-lesson layout in this deck")
    bullets(s, MARGIN, 1.7, SLIDE_W - 2 * MARGIN, 5.2, [
        "1️⃣  Divider — lesson title, Pokémon image",
        "2️⃣  Dataset — which CSVs, which joins",
        "3️⃣  Instructor demo — every sheet, in order",
        "4️⃣  Calculated fields — every formula",
        "5️⃣  Dashboard — size and sheet order (when present)",
        "6️⃣  Student assignment — blank path, what to build, answer-key path",
    ], size=20, spacing=16)


def agenda(prs, day):
    s = blank_slide(prs)
    header_bar(s, NAVY, f"Day {day} — Agenda")
    items = [(slug, m) for slug, m in LESSON.items() if m["day"] == day]
    rows = []
    for slug, m in items:
        num = slug.split("_")[0]
        rows.append(f"Lesson {num} — {m['title']}  ({m['duration']})")
        rows.append(f"  {m['subtitle']}")
    bullets(s, MARGIN, 1.7, SLIDE_W - 2 * MARGIN, 5.5, rows, size=18, spacing=6)


def section_divider(prs, slug):
    m = LESSON[slug]
    s = blank_slide(prs, RED)
    num = slug.split("_")[0]
    text(s, MARGIN, 1.2, SLIDE_W - 2 * MARGIN, 0.7,
         f"LESSON {num}", size=30, bold=True, color=YELLOW, align="center")
    text(s, MARGIN, 2.0, SLIDE_W - 2 * MARGIN, 1.4,
         m["title"], size=58, bold=True, color=WHITE, align="center")
    text(s, MARGIN, 3.6, SLIDE_W - 2 * MARGIN, 0.7,
         m["subtitle"], size=22, color=WHITE, align="center")
    text(s, MARGIN, 4.4, SLIDE_W - 2 * MARGIN, 0.5,
         f"{m['duration']} · Day {m['day']}", size=18, color=YELLOW, align="center")
    image(s, SLIDE_W / 2 - 0.9, 5.1, 1.8, m["image"])


def dataset_slide(prs, slug, cfg):
    m = LESSON[slug]
    s = blank_slide(prs)
    header_bar(s, NAVY, f"Lesson {slug.split('_')[0]} — Dataset")
    items = []
    if "tables" in cfg:
        items.append(f"{len(cfg['tables'])} CSVs joined inside Tableau:")
        for t in cfg["tables"]:
            items.append(f"  {Path(t['csv_path']).name}  (alias: {t['alias']})")
        if cfg.get("joins"):
            items.append("")
            items.append("Joins configured:")
            for j in cfg["joins"]:
                items.append(f"  {j['left_table']}.{j['left_col']} = {j['right_table']}.{j['right_col']}  ({j.get('join_type', 'inner')})")
    else:
        items.append(f"{Path(cfg['data_csv']).name}")
    # Image on the right
    bullets(s, MARGIN, 1.7, 8.2, 5.4, items, size=18, spacing=8)
    image(s, 9.5, 2.0, 3.2, m["image"])


def demo_flow_slides(prs, slug, cfg):
    """Max 4 sheets per slide. Each sheet gets a clean block."""
    sheets = cfg.get("worksheets", [])
    chunks = [sheets[i:i+4] for i in range(0, len(sheets), 4)] or [[]]
    total = len(chunks)
    for idx, chunk in enumerate(chunks, 1):
        s = blank_slide(prs)
        suffix = f"  ({idx}/{total})" if total > 1 else ""
        header_bar(s, RED, f"Lesson {slug.split('_')[0]} — Instructor Demo{suffix}")
        # Intro line
        text(s, MARGIN, 1.5, SLIDE_W - 2 * MARGIN, 0.4,
             "Walk these sheets in the workbook, in order:",
             size=16, color=GRAY)
        # Per-sheet blocks
        y = 2.05
        for i, ws in enumerate(chunk, 1 + (idx - 1) * 4):
            row_str = ", ".join(_fmt_field(f) for f in ws.get("rows", [])) or "—"
            col_str = ", ".join(_fmt_field(f) for f in ws.get("cols", [])) or "—"
            enc_parts = []
            for enc in ("color", "size", "shape"):
                if enc in ws:
                    enc_parts.append(f"{enc}: {_fmt_field(ws[enc])}")
            enc_str = "  |  ".join(enc_parts)
            # Sheet title
            text(s, MARGIN, y, SLIDE_W - 2 * MARGIN, 0.4,
                 f"{i}.  {ws['name']}",
                 size=18, bold=True, color=NAVY)
            # Detail line
            detail = f"Mark: {ws.get('mark', 'Automatic')}   ·   Rows: {row_str}   ·   Cols: {col_str}"
            if enc_str:
                detail += f"\n{enc_str}"
            text(s, MARGIN + 0.4, y + 0.45, SLIDE_W - 2 * MARGIN - 0.4, 0.7,
                 detail, size=13, color=DARK)
            y += 1.2


def calcs_slides(prs, slug, cfg):
    calcs = cfg.get("calculated_fields", [])
    if not calcs:
        return
    # Max 3 calcs per slide
    chunks = [calcs[i:i+3] for i in range(0, len(calcs), 3)]
    total = len(chunks)
    for idx, chunk in enumerate(chunks, 1):
        s = blank_slide(prs)
        suffix = f"  ({idx}/{total})" if total > 1 else ""
        header_bar(s, NAVY, f"Lesson {slug.split('_')[0]} — Calculated Fields{suffix}")
        y = 1.6
        for cf in chunk:
            text(s, MARGIN, y, SLIDE_W - 2 * MARGIN, 0.4,
                 f"{cf['caption']}   ({cf['datatype']})",
                 size=18, bold=True, color=NAVY)
            # Formula in monospace-ish feel (use smaller regular font since python-pptx
            # doesn't reliably set monospace across Google Slides)
            formula = cf["formula"].replace("\n", " · ")
            text(s, MARGIN + 0.3, y + 0.45, SLIDE_W - 2 * MARGIN - 0.3, 1.2,
                 formula, size=14, color=DARK)
            y += 1.85


def dashboard_slide(prs, slug, cfg):
    d = cfg.get("dashboard")
    if not d:
        return
    s = blank_slide(prs)
    header_bar(s, NAVY, f"Lesson {slug.split('_')[0]} — Dashboard")
    w, h = d.get("size", [1200, 800])
    text(s, MARGIN, 1.7, SLIDE_W - 2 * MARGIN, 0.5,
         f"Name: {d['name']}    ·    Size: {w} × {h}",
         size=18, bold=True, color=DARK)
    text(s, MARGIN, 2.4, SLIDE_W - 2 * MARGIN, 0.5,
         "Sheets (top → bottom):", size=16, color=GRAY)
    bullets(s, MARGIN, 2.9, SLIDE_W - 2 * MARGIN, 4.2,
            [f"{i}.  {sh}" for i, sh in enumerate(d["sheets"], 1)],
            size=16, spacing=6)


def student_slide(prs, slug, cfg):
    m = LESSON[slug]
    s = blank_slide(prs)
    header_bar(s, GREEN, f"Lesson {slug.split('_')[0]} — Student Assignment")
    if slug == "13_master_reference":
        bullets(s, MARGIN, 1.7, SLIDE_W - 2 * MARGIN, 5.3, [
            "No student activity.",
            "This is the instructor's apex reference — used during office hours.",
            "Students do not submit work for this lesson.",
        ], size=20, spacing=14)
        return
    n_sheets = len(cfg.get("worksheets", []))
    n_calcs = len(cfg.get("calculated_fields", []))
    items = [
        f"Time budget: {m['duration']}",
        "",
        f"1. Open the blank workbook:  student/workbooks_blank/{slug}_blank.twbx",
        f"2. Build {n_sheets} sheet{'s' if n_sheets != 1 else ''} matching the instructor demo",
    ]
    if n_calcs:
        items.append(f"3. Create {n_calcs} calculated field{'s' if n_calcs != 1 else ''} (formulas on the calc slide)")
    items.extend([
        "",
        f"Stuck? Compare against:  student/workbooks/{slug}_student.twbx",
        f"Save your file as:  {slug}_<your-initials>.twbx",
    ])
    bullets(s, MARGIN, 1.7, SLIDE_W - 2 * MARGIN, 5.3, items, size=18, spacing=10)


def closing(prs):
    s = blank_slide(prs, NAVY)
    text(s, MARGIN, 2.4, SLIDE_W - 2 * MARGIN, 1.3,
         "That's a wrap, Trainer.", size=54, bold=True, color=WHITE, align="center")
    text(s, MARGIN, 3.9, SLIDE_W - 2 * MARGIN, 0.7,
         "All artifacts live in course/ — commit and share.",
         size=22, color=YELLOW, align="center")
    image(s, SLIDE_W / 2 - 1.0, 4.9, 2.0, "pikachu.jpg")


def _fmt_field(f):
    agg = f.get("agg", "None")
    return f"[{f['field']}]" if agg == "None" else f"{agg}([{f['field']}])"


# ---------------- Build ----------------

def build():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    title_slide(prs)
    how_to_use(prs)
    per_lesson_layout(prs, 1)
    agenda(prs, 1)
    agenda(prs, 2)

    for slug in LESSON:
        cfg_path = CONFIGS_DIR / f"{slug}_instructor.json"
        if not cfg_path.exists():
            cfg_path = CONFIGS_DIR / f"{slug}.json"
        if not cfg_path.exists():
            print(f"missing config: {slug}", file=sys.stderr)
            continue
        cfg = json.loads(cfg_path.read_text())
        section_divider(prs, slug)
        dataset_slide(prs, slug, cfg)
        demo_flow_slides(prs, slug, cfg)
        calcs_slides(prs, slug, cfg)
        dashboard_slide(prs, slug, cfg)
        student_slide(prs, slug, cfg)

    closing(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"wrote {OUT} ({len(prs.slides)} slides, {OUT.stat().st_size} bytes)")


if __name__ == "__main__":
    build()
